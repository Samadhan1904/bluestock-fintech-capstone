# generate_presentation.py
# Bluestock Fintech Capstone
# Day 7 - Task 2: Generate the 12-slide presentation deck
#
# Builds Bluestock_MF_Presentation.pptx using python-pptx
# Pulls real numbers from our processed CSVs

import pandas as pd
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


BASE_DIR    = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
CLEAN_PATH  = os.path.join(BASE_DIR, "data", "processed")
CHART_PATH  = os.path.join(BASE_DIR, "reports", "charts")
REPORT_PATH = os.path.join(BASE_DIR, "reports")

OUTPUT_FILE = os.path.join(REPORT_PATH, "Bluestock_MF_Presentation.pptx")

# colors
BLUE  = RGBColor(0x19, 0x76, 0xD2)
DARK  = RGBColor(0x0D, 0x47, 0xA1)
GRAY  = RGBColor(0x60, 0x60, 0x60)

# load data we'll reference
scorecard = pd.read_csv(os.path.join(CLEAN_PATH, "fund_scorecard.csv"))


# create presentation with 16:9 widescreen size
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # blank layout


def add_title(slide, text, subtitle=None):
    """Add a title text box at the top of the slide."""
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BLUE

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = GRAY


def add_bullets(slide, bullets, top=1.8, left=0.8, width=11.7, height=5,
                font_size=18):
    """Add a bulleted text box."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {bullet}"
        p.font.size = Pt(font_size)
        p.space_after = Pt(10)


def add_image(slide, path, left, top, width=None, height=None):
    """Add an image if it exists, otherwise add a placeholder text."""
    if os.path.exists(path):
        if width:
            slide.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(width))
        else:
            slide.shapes.add_picture(path, Inches(left), Inches(top), height=Inches(height))
    else:
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(6), Inches(1))
        box.text_frame.text = f"[Image not found: {os.path.basename(path)}]"


def add_footer(slide, page_num):
    box = slide.shapes.add_textbox(Inches(11.8), Inches(7.05), Inches(1.3), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    p.text = f"Bluestock Fintech | {page_num}/12"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT


# ===================================================================
# SLIDE 1: TITLE
# ===================================================================
slide = prs.slides.add_slide(blank_layout)

box = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(11.3), Inches(2.5))
tf = box.text_frame
p = tf.paragraphs[0]
p.text = "Bluestock Fintech"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = BLUE
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "Mutual Fund Analytics Platform"
p2.font.size = Pt(28)
p2.font.color.rgb = DARK
p2.alignment = PP_ALIGN.CENTER

p3 = tf.add_paragraph()
p3.text = "Capstone Project Presentation"
p3.font.size = Pt(18)
p3.font.color.rgb = GRAY
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(20)

add_footer(slide, 1)


# ===================================================================
# SLIDE 2: PROBLEM & OBJECTIVE
# ===================================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Problem & Objective")
add_bullets(slide, [
    "India's mutual fund industry manages Rs.81 lakh crore across 1,908 schemes "
    "and 26.12 crore investor folios, but data is fragmented across AMFI, NSE/BSE "
    "and third-party APIs.",

    "Investors struggle to compare funds on a risk-adjusted basis - raw NAV data "
    "requires significant transformation to compute Sharpe, Alpha, Beta, etc.",

    "Most retail investors don't know if their fund is beating its benchmark.",

    "Objective: Build an end-to-end analytics platform - ETL pipeline, SQLite "
    "database, EDA, performance/risk metrics, interactive dashboard, and "
    "investor behaviour analytics.",
], font_size=20)
add_footer(slide, 2)


# ===================================================================
# SLIDE 3: DATA SOURCES
# ===================================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Data Sources & Datasets")
add_bullets(slide, [
    "10 provided CSV datasets: fund master (40 schemes), NAV history (46,000 "
    "rows), AUM by fund house, monthly SIP inflows, category inflows, folio "
    "counts, scheme performance, investor transactions (32,778 rows), portfolio "
    "holdings, benchmark indices (8,050 rows).",

    "Live data: mfapi.in REST API used to fetch real-time NAV for 6 large-cap "
    "schemes (HDFC Top 100, SBI/ICICI/Axis/Kotak Bluechip, Nippon Large Cap).",

    "All AMFI codes, fund names, expense ratios and AUM figures sourced from "
    "real AMFI India / mfapi.in data. Investor transactions are simulated using "
    "real demographic and geographic distributions.",
], font_size=20)
add_footer(slide, 3)


# ===================================================================
# SLIDE 4: ARCHITECTURE
# ===================================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "System Architecture")
add_bullets(slide, [
    "Layer 1 - Extract: 10 CSV files + mfapi.in REST API",
    "Layer 2 - Transform: Python/Pandas - date parsing, forward-fill, "
    "deduplication, validation, daily return calculation",
    "Layer 3 - Load: SQLite star schema - 1 dimension table (dim_fund) + "
    "6 fact tables (nav, transactions, performance, portfolio, AUM, SIP), "
    "79,318 total rows",
    "Layer 4 - Analyse: Jupyter notebooks - EDA, Sharpe/Sortino/Alpha/Beta/VaR/HHI",
    "Layer 5 - Visualise: Power BI dashboard - 4 interactive pages with slicers",
], font_size=20)
add_footer(slide, 4)


# ===================================================================
# SLIDE 5-6: EDA HIGHLIGHTS
# ===================================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "EDA Highlights (1/2)")
add_image(slide, os.path.join(CHART_PATH, "nav_trends_equity.png"),
          left=0.6, top=1.6, width=6.0)
add_image(slide, os.path.join(CHART_PATH, "sip_inflow_trend.png"),
          left=6.9, top=1.6, width=6.0)
add_bullets(slide, [
    "All 40 equity funds show consistent NAV growth 2022-2026",
    "Monthly SIP inflows grew 169% - Rs.11,517 Cr to Rs.31,002 Cr (Dec 2025 ATH)",
], top=5.7, font_size=16)
add_footer(slide, 5)

slide = prs.slides.add_slide(blank_layout)
add_title(slide, "EDA Highlights (2/2)")
add_image(slide, os.path.join(CHART_PATH, "folio_count_growth.png"),
          left=0.6, top=1.6, width=6.0)
add_image(slide, os.path.join(CHART_PATH, "sector_allocation_donut.png"),
          left=6.9, top=1.6, width=6.0)
add_bullets(slide, [
    "Folio count nearly doubled - 13.26 Cr to 26.12 Cr (97% growth in 4 years)",
    "Banking & Financial Services dominates sector allocation (25-30%)",
], top=5.7, font_size=16)
add_footer(slide, 6)


# ===================================================================
# SLIDE 7-8: PERFORMANCE METRICS
# ===================================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Performance Metrics (1/2)")
add_bullets(slide, [
    "Annualised returns computed for all 40 funds using 252 trading days",
    "CAGR (1yr, 3yr, 4yr) - Best 3yr: Axis Midcap at 36.07%",
    "Sharpe Ratio (Rf=6.5%) - Best: Mirae Asset Large Cap at 1.45",
    "Sortino Ratio - Best: Mirae Asset Large Cap at 2.39",
    "Alpha & Beta vs Nifty 100 computed via OLS regression",
    "Max Drawdown - Worst: SBI Small Cap Direct at -52.57% "
    "(needed 110.86% to recover)",
], font_size=20)
add_footer(slide, 7)

slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Performance Metrics (2/2) - Fund Scorecard")
add_image(slide, os.path.join(CHART_PATH, "fund_scorecard.png"),
          left=2.5, top=1.5, height=5.5)
add_footer(slide, 8)


# ===================================================================
# SLIDE 9-10: DASHBOARD SCREENSHOTS
# ===================================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Interactive Dashboard (1/2)")
add_image(slide, os.path.join(CHART_PATH, "dashboard_page1_industry.jpeg"),
          left=0.6, top=1.6, width=6.0)
add_image(slide, os.path.join(CHART_PATH, "dashboard_page2_performance.jpeg"),
          left=6.9, top=1.6, width=6.0)
add_bullets(slide, [
    "Page 1: Industry Overview - KPI cards for AUM, SIP, Folios, Schemes",
    "Page 2: Fund Performance - Return vs Risk scatter, scorecard table",
], top=5.7, font_size=16)
add_footer(slide, 9)

slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Interactive Dashboard (2/2)")
add_image(slide, os.path.join(CHART_PATH, "dashboard_page3_investor.jpeg"),
          left=0.6, top=1.6, width=6.0)
add_image(slide, os.path.join(CHART_PATH, "dashboard_page4_sip_trends.jpeg"),
          left=6.9, top=1.6, width=6.0)
add_bullets(slide, [
    "Page 3: Investor Analytics - State/age/city-tier breakdowns",
    "Page 4: SIP & Market Trends - Inflow trends, category heatmap",
], top=5.7, font_size=16)
add_footer(slide, 10)


# ===================================================================
# SLIDE 11: KEY FINDINGS
# ===================================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Key Findings")
add_bullets(slide, [
    "SBI Mutual Fund leads with Rs.12.5 lakh crore AUM (largest AMC)",
    "All 40 funds beat their 3-year benchmark - positive alpha across the board",
    "Small Cap funds carry highest risk (VaR/CVaR) - up to -5.8% worst-day loss",
    "Liquid funds are near risk-free - worst-day loss under 0.1%",
    "97.8% of frequent SIP investors show gaps >35 days - SIP discipline gap",
    "2025 cohort invests 23% more per SIP (Rs.13,505) vs 2024 cohort (Rs.10,997)",
    "All equity portfolios have HHI < 2500 - reasonable sector diversification",
    "26-35 age group has the most SIP investors; 56+ invests highest average amount",
], font_size=18)
add_footer(slide, 11)


# ===================================================================
# SLIDE 12: THANK YOU
# ===================================================================
slide = prs.slides.add_slide(blank_layout)

box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.3), Inches(2))
tf = box.text_frame
p = tf.paragraphs[0]
p.text = "Thank You"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = BLUE
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "Bluestock Fintech | Mutual Fund Analytics Capstone"
p2.font.size = Pt(18)
p2.font.color.rgb = GRAY
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(20)

p3 = tf.add_paragraph()
p3.text = "GitHub: github.com/Samadhan1904/bluestock-fintech-capstone"
p3.font.size = Pt(14)
p3.font.color.rgb = GRAY
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(10)

add_footer(slide, 12)


# ── Save ──────────────────────────────────────────────────────
prs.save(OUTPUT_FILE)
print(f"Presentation generated: {OUTPUT_FILE}")
print(f"Total slides: {len(prs.slides.__iter__().__self__._sldIdLst)}")
